"""Tools for creating Office documents (Word, Excel, PowerPoint, PDF)."""

import logging
import re
from datetime import datetime
from pathlib import Path

from langchain_core.tools import StructuredTool
from pydantic import BaseModel, Field


logger = logging.getLogger(__name__)


class CreateWordInput(BaseModel):
    path: str = Field(
        description=(
            "Ruta completa del archivo .docx. "
            "Ejemplo: C:/Users/User/Downloads/aiko_personal/guia.docx"
        )
    )
    title: str = Field(
        description=(
            "Título completo en español, mínimo 6 palabras. "
            "NUNCA una sola palabra como 'Dieta' o 'Guía'."
        )
    )
    content: str = Field(
        description=(
            "Contenido COMPLETO (mín. ~450 palabras) con secciones, "
            "viñetas (-) y pasos numerados (1. 2. 3.)."
        )
    )
    images: list[str] = Field(
        default_factory=list,
        description=(
            "Rutas locales de imágenes a insertar (jpg/png/webp). "
            "Ejemplo: ['C:/Users/User/Downloads/foto.jpg']. "
            "Si el usuario subió imágenes, usa esas rutas."
        ),
    )


class CreateExcelInput(BaseModel):
    path: str = Field(description="Ruta completa del archivo .xlsx")
    title: str = Field(description="Título del reporte")
    headers: list[str] = Field(description="Encabezados de columna")
    rows: list[list[str]] = Field(description="Filas de datos")
    chart_type: str = Field(default="bar", description="bar | pie | line | area")
    images: list[str] = Field(
        default_factory=list,
        description="Rutas de imágenes opcionales a colocar arriba de la tabla",
    )


class CreatePowerPointInput(BaseModel):
    path: str = Field(description="Ruta completa del archivo .pptx")
    title: str = Field(description="Título de la presentación")
    slides: list[dict] = Field(
        description="Lista de {title, content}. content puede ser str o lista"
    )
    theme: str = Field(default="auto", description="auto|tech|business|education|health|minimal")
    images: list[str] = Field(
        default_factory=list,
        description="Imágenes globales (portada / primera slide de contenido)",
    )


class CreatePDFInput(BaseModel):
    path: str = Field(description="Ruta completa del archivo .pdf")
    title: str = Field(
        description="Título completo en español, mínimo 6 palabras"
    )
    content: str = Field(
        description=(
            "Contenido COMPLETO con secciones y viñetas. "
            "Mínimo recomendado ~450 palabras."
        )
    )
    images: list[str] = Field(
        default_factory=list,
        description=(
            "Rutas locales de imágenes a insertar en el PDF. "
            "Si el usuario subió fotos, pásalas aquí."
        ),
    )


class DocumentCreatorTool:
    def __init__(self, allowed_paths: list[str] | None = None):
        self.allowed_paths = [
            Path(p).resolve() for p in (allowed_paths or ["./documents", "./uploads"])
        ]
        self.extra_roots = [
            Path("C:/Users/User/Downloads").resolve(),
            Path("C:/Users/User/Documents").resolve(),
            Path("C:/Users/User/OneDrive/Documentos").resolve(),
            Path("C:/Users/aleja/Downloads").resolve(),
            Path("C:/Users/aleja/Documents").resolve(),
            Path("C:/Users/aleja/OneDrive/Documentos").resolve(),
            Path("./data/uploads").resolve(),
            Path("./uploads").resolve(),
        ]
        logger.info("✅ DocumentCreatorTool initialized")

    def _is_allowed(self, path: Path) -> bool:
        try:
            resolved = path.resolve()
            for allowed in self.allowed_paths + self.extra_roots:
                if str(resolved).lower().startswith(str(allowed).lower()):
                    return True
            return False
        except Exception:
            return False

    def _resolve_images(self, images: list[str] | None) -> list[Path]:
        """Filtra rutas de imagen existentes y permitidas."""
        out: list[Path] = []
        for raw in images or []:
            try:
                p = Path(str(raw).strip().strip('"').strip("'"))
                if not p.exists() or not p.is_file():
                    logger.warning(f"Imagen no encontrada: {p}")
                    continue
                if p.suffix.lower() not in {".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp"}:
                    logger.warning(f"Formato de imagen no soportado: {p}")
                    continue
                # Permitir si está en carpetas allowed O si es un upload del proyecto
                if not self._is_allowed(p):
                    # aún permitir uploads del backend
                    if "upload" not in str(p).lower() and "aiko" not in str(p).lower():
                        logger.warning(f"Imagen fuera de rutas permitidas: {p}")
                        continue
                out.append(p.resolve())
            except Exception as e:
                logger.warning(f"No se pudo usar imagen {raw}: {e}")
        return out[:8]  # máximo 8 imágenes

    @staticmethod
    def _clean_bullet(text: str) -> str:
        text = text.strip()
        text = re.sub(r"^[-•*–—]\s*", "", text)
        text = re.sub(r"^\d+[.)]\s*", "", text)
        return text.strip()

    @staticmethod
    def _is_markdown_heading(text: str) -> tuple[bool, int, str]:
        m = re.match(r"^(#{1,3})\s+(.+)$", text.strip())
        if m:
            return True, len(m.group(1)), m.group(2).strip()
        return False, 0, text

    @staticmethod
    def _is_subtitle(text: str, title: str) -> bool:
        t = text.strip()
        if not t or t.lower() == title.lower():
            return False
        if t.startswith(("#", "-", "•", "*", "–", "—")):
            return False
        if re.match(r"^\d+[.)]\s+", t):
            return False
        if len(t) > 100:
            return False
        if t.endswith(".") and len(t) > 60:
            return False
        if t.count(" ") > 12:
            return False
        return True

    @staticmethod
    def _is_bullet(text: str) -> bool:
        t = text.strip()
        if t.startswith(("-", "•", "*", "–", "—")):
            return True
        if re.match(r"^\d+[.)]\s+\S", t):
            return True
        return False

    @staticmethod
    def _escape_pdf(text: str) -> str:
        if not text:
            return ""
        return (
            str(text)
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
        )

    @staticmethod
    def _word_count(text: str) -> int:
        return len(re.findall(r"\b\w+\b", text or ""))

    @staticmethod
    def _expand_title(title: str, content: str = "") -> str:
        t = (title or "").strip()
        low = t.lower()
        blob = f"{t} {content}".lower()
        if (
            len(t.split()) >= 6
            and low not in {"guía", "guia", "diet", "dieta", "informe", "documento"}
        ):
            return t
        if any(
            k in blob
            for k in [
                "dieta",
                "alimentación",
                "alimentacion",
                "perder peso",
                "nutrición",
                "nutricion",
            ]
        ):
            return "Guía de dieta saludable para perder peso de forma sostenible"
        if any(
            k in blob
            for k in ["senati", "formación profesional", "formacion profesional"]
        ):
            return "Guía general sobre formación profesional y SENATI"
        if any(
            k in blob for k in ["ia ", "inteligencia artificial", "machine learning"]
        ):
            return "Informe sobre inteligencia artificial y aplicaciones actuales"
        if len(t.split()) < 4:
            return f"Documento informativo: {t}" if t else "Documento generado por Aiko"
        return t

    def _diet_template(self) -> str:
        return """Introducción
Una dieta saludable no es un castigo ni una moda de pocos días: es un patrón de alimentación que puedes mantener en el tiempo. El objetivo no es solo bajar de peso, sino nutrir el cuerpo, sostener la energía diaria y reducir riesgos de salud a largo plazo.

Principios básicos
- Prioriza alimentos reales: frutas, verduras, legumbres, huevos, pescado, carnes magras, lácteos según tolerancia y cereales integrales.
- Incluye proteína en cada comida principal para mayor saciedad y para cuidar la masa muscular.
- Elige carbohidratos de calidad (avena, arroz, papa, quinoa, pan integral) según tu nivel de actividad.
- Usa grasas saludables (aceite de oliva, palta, frutos secos) y limita frituras habituales.
- Bebe agua de forma regular; reduce bebidas azucaradas y alcohol.
- Si buscas perder grasa, un déficit calórico moderado suele ser más sostenible que una restricción extrema.

Cómo armar el plato
- Mitad del plato: verduras o ensalada.
- Un cuarto: proteína (pollo, pescado, huevo, legumbres, tofu).
- Un cuarto: carbohidrato complejo.
- Añade una grasa saludable en cantidad controlada.

Ejemplos orientativos
- Desayuno: avena con fruta y yogur, o huevos con pan integral y verdura.
- Almuerzo: pollo o pescado a la plancha, arroz o papa, y ensalada abundante.
- Cena: similar al almuerzo, algo más ligera si cenas tarde.
- Snacks: fruta, yogur natural, un puñado de frutos secos o vegetales con hummus.

Pasos prácticos para empezar
1. Define un objetivo realista (hábitos sostenibles antes que kilos a toda costa).
2. Planifica 3 comidas principales y, si hace falta, 1 snack.
3. Cocina más en casa para controlar ingredientes y porciones.
4. Lee etiquetas: menos azúcar añadida y menos ultraprocesados.
5. Combina alimentación con caminata diaria o fuerza 2–3 veces por semana.
6. Prioriza el sueño: el cansancio aumenta antojos y baja la adherencia.

Errores comunes
- Déficit muy agresivo que agota y termina en atracones.
- Eliminar por completo carbohidratos o grasas sin criterio.
- Pesarse cada día y obsesionarse con la báscula.
- Confiar solo en productos dietéticos ultraprocesados.
- No ajustar el plan cuando cambian el trabajo, el estrés o el deporte.

Conclusión
La mejor dieta es la que puedes sostener: variada, suficiente en proteína, moderada en calorías si buscas perder grasa, y compatible con tu vida real. Si tienes una condición médica, embarazo o tomas medicación, consulta a un profesional de la salud antes de cambios importantes.
"""

    def _generic_template(self, title: str) -> str:
        return f"""Introducción
Este documento ofrece una visión clara y ordenada sobre el tema: {title}. El propósito es explicar ideas clave, pasos útiles y precauciones para aplicar el conocimiento de forma práctica.

Puntos clave
- Comprende el objetivo principal antes de entrar en detalles.
- Separa lo esencial de lo accesorio.
- Aplica los conceptos con ejemplos concretos.
- Revisa resultados y ajusta cuando haga falta.

Desarrollo
1. Contexto: por qué importa este tema.
2. Conceptos básicos que conviene manejar.
3. Pasos recomendados para ponerlo en práctica.
4. Recursos o apoyos útiles.
5. Cómo medir si estás avanzando.

Errores frecuentes
- Querer resultados inmediatos sin proceso.
- Copiar soluciones ajenas sin adaptarlas a tu caso.
- Ignorar límites de tiempo, presupuesto o salud.

Conclusión
Un enfoque ordenado reduce confusión y mejora la constancia. Este material es orientativo; complementa con fuentes confiables y, cuando corresponda, con asesoría profesional.
"""

    def _ensure_rich_content(self, title: str, content: str) -> tuple[str, str]:
        title = self._expand_title(title, content)
        text = (content or "").strip()
        words = self._word_count(text)
        low = f"{title} {text}".lower()

        has_sections = any(
            h.lower() in text.lower()
            for h in ["Introducción", "Principios", "Pasos", "Errores", "Conclusión"]
        )
        is_diet = any(
            k in low
            for k in [
                "dieta",
                "alimentación",
                "alimentacion",
                "nutri",
                "perder peso",
                "calorías",
                "calorias",
                "saludable",
            ]
        )

        if words >= 280 and has_sections:
            return title, text
        if is_diet and words < 280:
            return title, self._diet_template()
        if words < 200:
            return title, self._generic_template(title)
        if words < 280 and not has_sections:
            return title, f"{text}\n\n{self._generic_template(title)}"
        return title, text

    # ───────────────────── WORD ─────────────────────
    def create_word(
        self,
        path: str,
        title: str,
        content: str,
        images: list[str] | None = None,
    ) -> str:
        try:
            title, content = self._ensure_rich_content(title, content)
            if self._word_count(content) < 80:
                return "❌ Error: contenido demasiado corto."

            file_path = Path(path)
            if file_path.suffix.lower() != ".docx":
                file_path = file_path.with_suffix(".docx")
            if not self._is_allowed(file_path):
                return f"❌ Ruta no permitida: {file_path}"

            from docx import Document
            from docx.shared import Pt, Cm, RGBColor, Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.oxml.ns import qn
            from docx.oxml import OxmlElement

            doc = Document()
            for section in doc.sections:
                section.top_margin = Cm(2.0)
                section.bottom_margin = Cm(2.0)
                section.left_margin = Cm(2.3)
                section.right_margin = Cm(2.3)

            try:
                core = doc.core_properties
                core.author = "Aiko AI"
                core.title = title
            except Exception:
                pass

            # Barra superior simulada
            bar = doc.add_paragraph()
            bar.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = bar.add_run("━" * 28)
            run.font.color.rgb = RGBColor(37, 99, 235)
            run.font.size = Pt(14)

            heading = doc.add_heading(title.strip(), level=0)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in heading.runs:
                run.font.color.rgb = RGBColor(30, 64, 175)
                run.font.size = Pt(22)
                run.font.name = "Calibri"

            subtitle = doc.add_paragraph()
            subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = subtitle.add_run(
                f"Documento generado por Aiko · {datetime.now().strftime('%d/%m/%Y %H:%M')}"
            )
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(100, 116, 139)
            run.italic = True

            sep = doc.add_paragraph("─" * 42)
            sep.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in sep.runs:
                run.font.color.rgb = RGBColor(203, 213, 225)

            # Imágenes al inicio
            img_paths = self._resolve_images(images)
            for img in img_paths:
                try:
                    p = doc.add_paragraph()
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    run = p.add_run()
                    run.add_picture(str(img), width=Inches(5.2))
                    cap = doc.add_paragraph()
                    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    r = cap.add_run(f"Figura: {img.name}")
                    r.font.size = Pt(9)
                    r.font.color.rgb = RGBColor(100, 116, 139)
                    r.italic = True
                except Exception as e:
                    logger.warning(f"No se pudo insertar imagen en Word: {e}")

            raw = content.replace("\r\n", "\n").replace("\r", "\n")
            paragraphs = [p.strip() for p in raw.split("\n") if p.strip()]

            for para_text in paragraphs:
                if para_text.lower() == title.strip().lower():
                    continue
                is_md, level, md_text = self._is_markdown_heading(para_text)
                if is_md:
                    h = doc.add_heading(md_text, level=min(level, 3))
                    for run in h.runs:
                        run.font.color.rgb = RGBColor(37, 99, 235)
                    continue
                if self._is_bullet(para_text):
                    clean = self._clean_bullet(para_text)
                    if not clean:
                        continue
                    if re.match(r"^\d+[.)]\s+", para_text.strip()):
                        p = doc.add_paragraph(clean, style="List Number")
                    else:
                        p = doc.add_paragraph(clean, style="List Bullet")
                    for run in p.runs:
                        run.font.size = Pt(11)
                        run.font.name = "Calibri"
                        run.font.color.rgb = RGBColor(30, 41, 59)
                    continue
                if self._is_subtitle(para_text, title):
                    h = doc.add_heading(para_text, level=1)
                    for run in h.runs:
                        run.font.color.rgb = RGBColor(37, 99, 235)
                        run.font.size = Pt(14)
                    continue
                p = doc.add_paragraph(para_text)
                p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                p.paragraph_format.space_after = Pt(10)
                p.paragraph_format.line_spacing = 1.15
                for run in p.runs:
                    run.font.size = Pt(11)
                    run.font.name = "Calibri"
                    run.font.color.rgb = RGBColor(30, 41, 59)

            footer = doc.add_paragraph()
            footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = footer.add_run("— Generado automáticamente por Aiko AI —")
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(148, 163, 184)
            run.italic = True

            file_path.parent.mkdir(parents=True, exist_ok=True)
            doc.save(str(file_path))
            size_kb = file_path.stat().st_size / 1024
            logger.info(f"✅ Word creado: {file_path} ({size_kb:.1f} KB)")
            extra = f"\n🖼️ Imágenes: {len(img_paths)}" if img_paths else ""
            return (
                f"✅ Documento Word creado correctamente\n"
                f"📄 Archivo: {file_path}\n"
                f"📝 Título: {title}\n"
                f"📦 Tamaño: {size_kb:.1f} KB{extra}"
            )
        except Exception as e:
            logger.error(f"Error creating Word: {e}", exc_info=True)
            return f"❌ Error al crear Word: {str(e)}"

    # ───────────────────── EXCEL ─────────────────────
    def create_excel(
        self,
        path: str,
        title: str,
        headers: list[str],
        rows: list[list[str]],
        chart_type: str = "bar",
        images: list[str] | None = None,
    ) -> str:
        try:
            if not headers:
                return "❌ Error: falta headers."
            if rows is None:
                rows = []

            file_path = Path(path)
            if file_path.suffix.lower() != ".xlsx":
                file_path = file_path.with_suffix(".xlsx")
            if not self._is_allowed(file_path):
                return f"❌ Ruta no permitida: {file_path}"

            from openpyxl import Workbook
            from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
            from openpyxl.utils import get_column_letter
            from openpyxl.chart import BarChart, PieChart, LineChart, AreaChart, Reference
            from openpyxl.chart.label import DataLabelList
            from openpyxl.drawing.image import Image as XLImage
            from openpyxl.worksheet.table import Table, TableStyleInfo

            wb = Workbook()
            ws = wb.active
            ws.title = (title or "Hoja1").strip()[:31] or "Hoja1"

            title_font = Font(name="Calibri", size=16, bold=True, color="1E40AF")
            header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
            header_fill = PatternFill("solid", fgColor="2563EB")
            cell_font = Font(name="Calibri", size=11, color="1E293B")
            thin = Border(
                left=Side(style="thin", color="E2E8F0"),
                right=Side(style="thin", color="E2E8F0"),
                top=Side(style="thin", color="E2E8F0"),
                bottom=Side(style="thin", color="E2E8F0"),
            )
            center = Alignment(horizontal="center", vertical="center", wrap_text=True)
            left = Alignment(horizontal="left", vertical="center", wrap_text=True)
            right = Alignment(horizontal="right", vertical="center")

            col_count = max(len(headers), 1)
            start_row = 1

            # Imagen arriba
            img_paths = self._resolve_images(images)
            if img_paths:
                try:
                    xl_img = XLImage(str(img_paths[0]))
                    xl_img.width = 320
                    xl_img.height = 180
                    ws.add_image(xl_img, "A1")
                    ws.row_dimensions[1].height = 140
                    start_row = 10
                except Exception as e:
                    logger.warning(f"No se pudo insertar imagen en Excel: {e}")
                    start_row = 1

            ws.merge_cells(
                start_row=start_row,
                start_column=1,
                end_row=start_row,
                end_column=col_count,
            )
            title_cell = ws.cell(row=start_row, column=1, value=title)
            title_cell.font = title_font
            title_cell.alignment = Alignment(horizontal="center", vertical="center")
            ws.row_dimensions[start_row].height = 28

            ws.merge_cells(
                start_row=start_row + 1,
                start_column=1,
                end_row=start_row + 1,
                end_column=col_count,
            )
            date_cell = ws.cell(
                row=start_row + 1,
                column=1,
                value=f"Generado por Aiko · {datetime.now().strftime('%d/%m/%Y %H:%M')}",
            )
            date_cell.font = Font(name="Calibri", size=9, italic=True, color="64748B")
            date_cell.alignment = Alignment(horizontal="center")

            header_row = start_row + 3
            for col_idx, header in enumerate(headers, start=1):
                cell = ws.cell(row=header_row, column=col_idx, value=str(header))
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = center
                cell.border = thin

            def to_number(value):
                if isinstance(value, (int, float)):
                    return value
                if value is None:
                    return ""
                s = str(value).replace("S/", "").replace(",", "").strip()
                try:
                    if s == "":
                        return ""
                    return float(s) if "." in s else int(s)
                except Exception:
                    return value

            for row_idx, row in enumerate(rows, start=header_row + 1):
                if not isinstance(row, (list, tuple)):
                    row = [row]
                for col_idx in range(1, len(headers) + 1):
                    raw_val = row[col_idx - 1] if col_idx - 1 < len(row) else ""
                    header_lower = str(headers[col_idx - 1]).lower()
                    numeric_keys = [
                        "precio",
                        "igv",
                        "total",
                        "monto",
                        "venta",
                        "cantidad",
                        "cant",
                        "stock",
                    ]
                    cell_value = (
                        to_number(raw_val)
                        if any(k in header_lower for k in numeric_keys)
                        else raw_val
                    )
                    cell = ws.cell(row=row_idx, column=col_idx, value=cell_value)
                    cell.font = cell_font
                    cell.border = thin
                    if any(k in header_lower for k in ["precio", "total", "monto", "venta"]):
                        cell.alignment = right
                        if isinstance(cell_value, (int, float)):
                            cell.number_format = "#,##0.00"
                    elif any(k in header_lower for k in ["cantidad", "cant", "stock"]):
                        cell.alignment = center
                    else:
                        cell.alignment = left
                    if (row_idx - header_row) % 2 == 0:
                        cell.fill = PatternFill("solid", fgColor="F8FAFC")

            last_data_row = header_row + max(len(rows), 1)
            for col_idx, header in enumerate(headers, start=1):
                max_len = len(str(header))
                for row in rows:
                    if isinstance(row, (list, tuple)) and col_idx - 1 < len(row):
                        max_len = max(max_len, len(str(row[col_idx - 1])))
                ws.column_dimensions[get_column_letter(col_idx)].width = min(
                    max(max_len + 4, 12), 40
                )

            if len(rows) > 0:
                table_ref = (
                    f"A{header_row}:{get_column_letter(len(headers))}{last_data_row}"
                )
                try:
                    table = Table(displayName="TablaAiko", ref=table_ref)
                    table.tableStyleInfo = TableStyleInfo(
                        name="TableStyleMedium2",
                        showRowStripes=True,
                    )
                    ws.add_table(table)
                except Exception:
                    pass

            ct = (chart_type or "bar").strip().lower()
            if ct in ["circular", "pastel", "pie"]:
                ct = "pie"
            elif ct in ["linea", "línea", "line"]:
                ct = "line"
            elif ct in ["area", "área"]:
                ct = "area"
            else:
                ct = "bar"

            if len(rows) > 0 and len(headers) >= 2:
                data = Reference(
                    ws, min_col=len(headers), min_row=header_row, max_row=last_data_row
                )
                cats = Reference(
                    ws, min_col=1, min_row=header_row + 1, max_row=last_data_row
                )
                if ct == "pie":
                    chart = PieChart()
                    chart.add_data(data, titles_from_data=True)
                    chart.set_categories(cats)
                    chart.dataLabels = DataLabelList()
                    chart.dataLabels.showPercent = True
                elif ct == "line":
                    chart = LineChart()
                    chart.add_data(data, titles_from_data=True)
                    chart.set_categories(cats)
                elif ct == "area":
                    chart = AreaChart()
                    chart.add_data(data, titles_from_data=True)
                    chart.set_categories(cats)
                else:
                    chart = BarChart()
                    chart.type = "col"
                    chart.grouping = "clustered"
                    chart.add_data(data, titles_from_data=True)
                    chart.set_categories(cats)
                chart.title = title[:50] if title else "Resumen"
                chart.width = 15
                chart.height = 10
                ws.add_chart(chart, f"A{last_data_row + 3}")

            file_path.parent.mkdir(parents=True, exist_ok=True)
            wb.save(str(file_path))
            size_kb = file_path.stat().st_size / 1024
            logger.info(f"✅ Excel creado: {file_path} ({size_kb:.1f} KB)")
            return (
                f"✅ Excel creado correctamente\n"
                f"📊 Archivo: {file_path}\n"
                f"📝 Título: {title}\n"
                f"📦 Tamaño: {size_kb:.1f} KB"
            )
        except Exception as e:
            logger.error(f"Error creating Excel: {e}", exc_info=True)
            return f"❌ Error al crear Excel: {str(e)}"

    # ───────────────────── POWERPOINT ─────────────────────
    def create_powerpoint(
        self,
        path: str,
        title: str,
        slides: list[dict],
        theme: str = "auto",
        images: list[str] | None = None,
    ) -> str:
        try:
            if not slides:
                return "❌ Error: indica al menos una diapositiva."

            file_path = Path(path)
            if file_path.suffix.lower() != ".pptx":
                file_path = file_path.with_suffix(".pptx")
            if not self._is_allowed(file_path):
                return f"❌ Ruta no permitida: {file_path}"

            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.enum.text import PP_ALIGN

            themes = {
                "tech": {
                    "name": "Tech",
                    "primary": RGBColor(37, 99, 235),
                    "cover_bg": RGBColor(15, 23, 42),
                    "card": RGBColor(241, 245, 249),
                    "dark": RGBColor(15, 23, 42),
                    "text": RGBColor(30, 41, 59),
                    "muted": RGBColor(148, 163, 184),
                    "white": RGBColor(255, 255, 255),
                },
                "business": {
                    "name": "Business",
                    "primary": RGBColor(30, 64, 175),
                    "cover_bg": RGBColor(15, 23, 42),
                    "card": RGBColor(248, 250, 252),
                    "dark": RGBColor(15, 23, 42),
                    "text": RGBColor(30, 41, 59),
                    "muted": RGBColor(100, 116, 139),
                    "white": RGBColor(255, 255, 255),
                },
                "health": {
                    "name": "Health",
                    "primary": RGBColor(5, 150, 105),
                    "cover_bg": RGBColor(6, 78, 59),
                    "card": RGBColor(236, 253, 245),
                    "dark": RGBColor(6, 78, 59),
                    "text": RGBColor(30, 41, 59),
                    "muted": RGBColor(100, 116, 139),
                    "white": RGBColor(255, 255, 255),
                },
                "education": {
                    "name": "Education",
                    "primary": RGBColor(8, 145, 178),
                    "cover_bg": RGBColor(8, 47, 73),
                    "card": RGBColor(240, 249, 255),
                    "dark": RGBColor(12, 74, 110),
                    "text": RGBColor(30, 41, 59),
                    "muted": RGBColor(100, 116, 139),
                    "white": RGBColor(255, 255, 255),
                },
                "minimal": {
                    "name": "Minimal",
                    "primary": RGBColor(15, 23, 42),
                    "cover_bg": RGBColor(15, 23, 42),
                    "card": RGBColor(248, 250, 252),
                    "dark": RGBColor(15, 23, 42),
                    "text": RGBColor(51, 65, 85),
                    "muted": RGBColor(148, 163, 184),
                    "white": RGBColor(255, 255, 255),
                },
            }
            chosen = (theme or "auto").strip().lower()
            if chosen == "auto":
                blob = f"{title} {slides}".lower()
                if any(k in blob for k in ["salud", "dieta", "nutri"]):
                    chosen = "health"
                elif any(k in blob for k in ["edu", "senati", "curso"]):
                    chosen = "education"
                elif any(k in blob for k in ["tech", "ia", "software"]):
                    chosen = "tech"
                else:
                    chosen = "business"
            if chosen not in themes:
                chosen = "business"
            T = themes[chosen]

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            blank = prs.slide_layouts[6]

            def add_rect(slide, left, top, width, height, color):
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, left, top, width, height
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = color
                shape.line.fill.background()
                return shape

            def normalize_bullets(content):
                if isinstance(content, list):
                    return [str(x).strip() for x in content if str(x).strip()][:6]
                raw = str(content or "").replace("\r\n", "\n")
                out = []
                for line in raw.split("\n"):
                    line = line.strip()
                    if not line:
                        continue
                    line = re.sub(r"^[-•*–—]\s*", "", line)
                    line = re.sub(r"^\d+[.)]\s*", "", line)
                    if line:
                        out.append(line)
                return out[:6] or ["(Sin contenido)"]

            # Portada
            cover = prs.slides.add_slide(blank)
            add_rect(cover, Inches(0), Inches(0), prs.slide_width, prs.slide_height, T["cover_bg"])
            add_rect(cover, Inches(0), Inches(0), Inches(0.28), prs.slide_height, T["primary"])
            tbox = cover.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11), Inches(2))
            tf = tbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = (title or "Presentación").strip()
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = T["white"]
            p.font.name = "Calibri"

            img_paths = self._resolve_images(images)
            if img_paths:
                try:
                    cover.shapes.add_picture(
                        str(img_paths[0]), Inches(8.5), Inches(4.2), width=Inches(4.0)
                    )
                except Exception as e:
                    logger.warning(f"Imagen portada PPT: {e}")

            sbox = cover.shapes.add_textbox(Inches(0.9), Inches(5.2), Inches(8), Inches(0.5))
            tf = sbox.text_frame
            p = tf.paragraphs[0]
            p.text = f"Generado por Aiko  ·  {datetime.now().strftime('%d/%m/%Y %H:%M')}"
            p.font.size = Pt(14)
            p.font.color.rgb = T["muted"]
            p.font.italic = True

            total = len(slides)
            for idx, item in enumerate(slides):
                if not isinstance(item, dict):
                    continue
                slide_title = str(item.get("title") or f"Diapositiva {idx + 1}").strip()
                bullets = normalize_bullets(item.get("content", ""))
                slide = prs.slides.add_slide(blank)
                add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, T["white"])
                add_rect(slide, Inches(0), Inches(0), Inches(0.18), prs.slide_height, T["primary"])
                add_rect(slide, Inches(0.18), Inches(0.1), Inches(13.15), Inches(1.15), T["card"])

                tbox = slide.shapes.add_textbox(Inches(0.55), Inches(0.4), Inches(12), Inches(0.6))
                tf = tbox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_title
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = T["dark"]

                for i, bullet in enumerate(bullets):
                    top = 1.55 + i * 0.8
                    cbox = slide.shapes.add_textbox(
                        Inches(0.9), Inches(top), Inches(11.5), Inches(0.55)
                    )
                    tf = cbox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = f"•  {bullet}"
                    p.font.size = Pt(16)
                    p.font.color.rgb = T["text"]

                # segunda imagen en la primera slide de contenido
                if idx == 0 and len(img_paths) > 1:
                    try:
                        slide.shapes.add_picture(
                            str(img_paths[1]),
                            Inches(9.8),
                            Inches(5.2),
                            width=Inches(3.0),
                        )
                    except Exception:
                        pass

                foot = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(6), Inches(0.3))
                p = foot.text_frame.paragraphs[0]
                p.text = f"Aiko AI  ·  {T['name']}"
                p.font.size = Pt(10)
                p.font.color.rgb = T["muted"]

                foot_r = slide.shapes.add_textbox(Inches(9.5), Inches(7.05), Inches(3.2), Inches(0.3))
                p = foot_r.text_frame.paragraphs[0]
                p.text = f"{idx + 1} / {total}"
                p.font.size = Pt(10)
                p.font.color.rgb = T["muted"]
                p.alignment = PP_ALIGN.RIGHT

            file_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(file_path))
            size_kb = file_path.stat().st_size / 1024
            logger.info(f"✅ PowerPoint creado: {file_path} ({size_kb:.1f} KB)")
            return (
                f"✅ PowerPoint creado\n"
                f"📊 Archivo: {file_path}\n"
                f"📝 Título: {title}\n"
                f"🎨 Tema: {T['name']}\n"
                f"📑 Diapositivas: {len(slides) + 1}\n"
                f"📦 Tamaño: {size_kb:.1f} KB"
            )
        except Exception as e:
            logger.error(f"Error creating PowerPoint: {e}", exc_info=True)
            return f"❌ Error al crear PowerPoint: {str(e)}"

    # ───────────────────── PDF ─────────────────────
    def create_pdf(
        self,
        path: str,
        title: str,
        content: str,
        images: list[str] | None = None,
    ) -> str:
        try:
            title, content = self._ensure_rich_content(title, content)
            if self._word_count(content) < 80:
                return "❌ Error: contenido demasiado corto."

            file_path = Path(path)
            if file_path.suffix.lower() != ".pdf":
                file_path = file_path.with_suffix(".pdf")
            if not self._is_allowed(file_path):
                return f"❌ Ruta no permitida: {file_path}"

            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import cm
            from reportlab.lib.colors import HexColor, white
            from reportlab.platypus import (
                SimpleDocTemplate,
                Paragraph,
                Spacer,
                HRFlowable,
                Image as RLImage,
                KeepTogether,
                Table,
                TableStyle,
            )
            from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT

            PRIMARY = HexColor("#2563EB")
            PRIMARY_DARK = HexColor("#1E40AF")
            TEXT = HexColor("#1E293B")
            MUTED = HexColor("#64748B")
            LINE = HexColor("#E2E8F0")
            BG_SOFT = HexColor("#EFF6FF")

            file_path.parent.mkdir(parents=True, exist_ok=True)

            doc = SimpleDocTemplate(
                str(file_path),
                pagesize=A4,
                rightMargin=2.0 * cm,
                leftMargin=2.0 * cm,
                topMargin=1.8 * cm,
                bottomMargin=1.8 * cm,
                title=title,
                author="Aiko AI",
            )

            styles = getSampleStyleSheet()
            style_title = ParagraphStyle(
                "AikoTitle",
                parent=styles["Title"],
                fontName="Helvetica-Bold",
                fontSize=20,
                textColor=PRIMARY_DARK,
                alignment=TA_CENTER,
                spaceAfter=4,
                leading=24,
            )
            style_meta = ParagraphStyle(
                "AikoMeta",
                parent=styles["Normal"],
                fontName="Helvetica-Oblique",
                fontSize=9,
                textColor=MUTED,
                alignment=TA_CENTER,
                spaceAfter=10,
            )
            style_h1 = ParagraphStyle(
                "AikoH1",
                parent=styles["Heading1"],
                fontName="Helvetica-Bold",
                fontSize=13,
                textColor=PRIMARY,
                spaceBefore=14,
                spaceAfter=6,
                leading=17,
            )
            style_body = ParagraphStyle(
                "AikoBody",
                parent=styles["Normal"],
                fontName="Helvetica",
                fontSize=11,
                textColor=TEXT,
                alignment=TA_JUSTIFY,
                spaceAfter=8,
                leading=15,
            )
            style_bullet = ParagraphStyle(
                "AikoBullet",
                parent=styles["Normal"],
                fontName="Helvetica",
                fontSize=11,
                textColor=TEXT,
                leftIndent=14,
                spaceAfter=4,
                leading=15,
            )
            style_caption = ParagraphStyle(
                "AikoCaption",
                parent=styles["Normal"],
                fontName="Helvetica-Oblique",
                fontSize=8,
                textColor=MUTED,
                alignment=TA_CENTER,
                spaceAfter=12,
            )
            style_footer = ParagraphStyle(
                "AikoFooter",
                parent=styles["Normal"],
                fontName="Helvetica-Oblique",
                fontSize=8,
                textColor=MUTED,
                alignment=TA_CENTER,
                spaceBefore=16,
            )

            story = []
            story.append(Spacer(1, 0.4 * cm))
            story.append(Paragraph(self._escape_pdf(title.strip()), style_title))
            story.append(
                Paragraph(
                    f"Documento generado por Aiko · {datetime.now().strftime('%d/%m/%Y %H:%M')}",
                    style_meta,
                )
            )
            story.append(
                HRFlowable(
                    width="100%",
                    thickness=2,
                    color=PRIMARY,
                    spaceBefore=2,
                    spaceAfter=12,
                )
            )

            # Imágenes
            img_paths = self._resolve_images(images)
            page_w = A4[0] - 4.0 * cm
            for img in img_paths:
                try:
                    im = RLImage(str(img))
                    # escalar al ancho útil
                    max_w = page_w * 0.92
                    max_h = 9 * cm
                    iw, ih = im.imageWidth, im.imageHeight
                    ratio = min(max_w / iw, max_h / ih, 1.0)
                    im.drawWidth = iw * ratio
                    im.drawHeight = ih * ratio
                    story.append(Spacer(1, 0.2 * cm))
                    story.append(im)
                    story.append(
                        Paragraph(
                            f"Figura: {self._escape_pdf(img.name)}",
                            style_caption,
                        )
                    )
                except Exception as e:
                    logger.warning(f"No se pudo insertar imagen en PDF: {e}")

            raw = content.replace("\r\n", "\n").replace("\r", "\n")
            paragraphs = [p.strip() for p in raw.split("\n") if p.strip()]

            for para_text in paragraphs:
                if para_text.lower() == title.strip().lower():
                    continue
                is_md, level, md_text = self._is_markdown_heading(para_text)
                if is_md:
                    story.append(Paragraph(self._escape_pdf(md_text), style_h1))
                    continue
                if self._is_bullet(para_text):
                    clean = self._clean_bullet(para_text)
                    if not clean:
                        continue
                    mnum = re.match(r"^(\d+)[.)]\s+", para_text.strip())
                    if mnum:
                        story.append(
                            Paragraph(
                                f"{mnum.group(1)}.  {self._escape_pdf(clean)}",
                                style_bullet,
                            )
                        )
                    else:
                        story.append(
                            Paragraph(f"•  {self._escape_pdf(clean)}", style_bullet)
                        )
                    continue
                if self._is_subtitle(para_text, title):
                    story.append(Paragraph(self._escape_pdf(para_text), style_h1))
                    continue
                story.append(Paragraph(self._escape_pdf(para_text), style_body))

            story.append(Spacer(1, 0.5 * cm))
            story.append(
                HRFlowable(
                    width="100%", thickness=0.8, color=LINE, spaceBefore=6, spaceAfter=6
                )
            )
            story.append(
                Paragraph("— Generado automáticamente por Aiko AI —", style_footer)
            )

            def _header_footer(canvas, doc_):
                canvas.saveState()
                canvas.setFillColor(PRIMARY)
                canvas.rect(0, A4[1] - 10, A4[0], 10, fill=1, stroke=0)
                canvas.setFillColor(MUTED)
                canvas.setFont("Helvetica", 8)
                canvas.drawString(2.0 * cm, 1.1 * cm, "Aiko AI Assistant")
                canvas.drawRightString(A4[0] - 2.0 * cm, 1.1 * cm, f"Pág. {doc_.page}")
                canvas.restoreState()

            doc.build(story, onFirstPage=_header_footer, onLaterPages=_header_footer)

            size_kb = file_path.stat().st_size / 1024
            logger.info(f"✅ PDF creado: {file_path} ({size_kb:.1f} KB)")
            extra = f"\n🖼️ Imágenes: {len(img_paths)}" if img_paths else ""
            return (
                f"✅ Documento PDF creado correctamente\n"
                f"📄 Archivo: {file_path}\n"
                f"📝 Título: {title}\n"
                f"📦 Tamaño: {size_kb:.1f} KB{extra}"
            )
        except Exception as e:
            logger.error(f"Error creating PDF: {e}", exc_info=True)
            return f"❌ Error al crear PDF: {str(e)}"

    def get_tools(self) -> list:
        return [
            StructuredTool.from_function(
                func=self.create_word,
                name="create_word",
                description=(
                    "Crea Word (.docx). "
                    "Puedes pasar images=[rutas] si el usuario subió fotos. "
                    "title ≥6 palabras; content largo con secciones."
                ),
                args_schema=CreateWordInput,
            ),
            StructuredTool.from_function(
                func=self.create_excel,
                name="create_excel",
                description=(
                    "Crea Excel (.xlsx) con tabla y gráfico. "
                    "images opcional (ruta de imagen arriba de la tabla)."
                ),
                args_schema=CreateExcelInput,
            ),
            StructuredTool.from_function(
                func=self.create_powerpoint,
                name="create_powerpoint",
                description=(
                    "Crea PowerPoint (.pptx). "
                    "images opcional para portada. Mínimo 3 slides."
                ),
                args_schema=CreatePowerPointInput,
            ),
            StructuredTool.from_function(
                func=self.create_pdf,
                name="create_pdf",
                description=(
                    "Crea PDF profesional. "
                    "Si el usuario subió imágenes, pásalas en images=[ruta1, ruta2]. "
                    "title ≥6 palabras; content con secciones."
                ),
                args_schema=CreatePDFInput,
            ),
        ]